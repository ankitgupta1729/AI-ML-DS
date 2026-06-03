"""Document loading, chunking and indexing.

Supports PDF, DOCX, PPTX, TXT and Markdown out of the box with a light
dependency footprint (no heavy ``unstructured`` install required).

For large corpora (e.g. multi-hundred-page previous-year-question PDFs) the
indexer keeps a small JSON *manifest* so unchanged files are skipped on the
next run — you only pay to embed what actually changed.
"""

from __future__ import annotations

import hashlib
import json
import logging
import time
from pathlib import Path
from typing import Iterable

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from .config import Settings, get_settings

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown"}

MANIFEST_NAME = "ingest_manifest.json"


# --------------------------------------------------------------------------- #
# File-type specific text extraction                                          #
# --------------------------------------------------------------------------- #
def _load_pdf(path: Path) -> list[Document]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    docs: list[Document] = []
    for i, page in enumerate(reader.pages):
        try:
            text = (page.extract_text() or "").strip()
        except Exception as exc:  # noqa: BLE001 - skip unreadable pages
            logger.debug("Page %d of %s unreadable: %s", i + 1, path.name, exc)
            continue
        if text:
            docs.append(Document(page_content=text, metadata={"page": i + 1}))
    return docs


def _load_docx(path: Path) -> list[Document]:
    import docx2txt

    text = (docx2txt.process(str(path)) or "").strip()
    return [Document(page_content=text)] if text else []


def _load_pptx(path: Path) -> list[Document]:
    from pptx import Presentation

    prs = Presentation(str(path))
    docs: list[Document] = []
    for i, slide in enumerate(prs.slides):
        parts = [
            shape.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        if parts:
            docs.append(
                Document(
                    page_content="\n".join(parts),
                    metadata={"slide": i + 1},
                )
            )
    return docs


def _load_text(path: Path) -> list[Document]:
    text = path.read_text(encoding="utf-8", errors="ignore").strip()
    return [Document(page_content=text)] if text else []


_LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
    ".txt": _load_text,
    ".md": _load_text,
    ".markdown": _load_text,
}


def load_file(path: Path) -> list[Document]:
    """Load a single supported file into a list of :class:`Document`."""
    loader = _LOADERS.get(path.suffix.lower())
    if loader is None:
        return []
    try:
        docs = loader(path)
    except Exception as exc:  # noqa: BLE001 - keep ingestion resilient
        logger.warning("Failed to load %s: %s", path, exc)
        return []

    # Tag every document with provenance metadata used for citations.
    for d in docs:
        d.metadata.setdefault("source", path.name)
        d.metadata["path"] = str(path)
        # Subject = the top-level folder under data/ (e.g. computer_science).
        d.metadata.setdefault("subject", _infer_subject(path))
    return docs


def _infer_subject(path: Path) -> str:
    settings = get_settings()
    try:
        rel = path.resolve().relative_to(settings.data_dir.resolve())
        return rel.parts[0] if len(rel.parts) > 1 else "general"
    except ValueError:
        return "general"


def discover_files(data_dir: Path) -> list[Path]:
    """Recursively find every supported document under ``data_dir``."""
    return sorted(
        p
        for p in data_dir.rglob("*")
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS
    )


def _chunk(documents: Iterable[Document], settings: Settings) -> list[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    chunks = splitter.split_documents(list(documents))
    # Stable, content-based IDs so re-ingesting the same material updates
    # existing vectors instead of creating duplicates.
    for c in chunks:
        key = f"{c.metadata.get('path', '')}:{c.page_content}"
        c.metadata["chunk_id"] = hashlib.sha256(key.encode()).hexdigest()[:16]
    return chunks


def _dedupe_by_id(chunks: list[Document]) -> list[Document]:
    """Drop chunks sharing a ``chunk_id`` (identical text repeated within a
    file — e.g. boilerplate headers/footers in PYQ PDFs).

    Chroma rejects duplicate IDs inside a single upsert request, and storing the
    same content twice adds no value, so we keep only the first occurrence.
    """
    seen: set[str] = set()
    unique: list[Document] = []
    for c in chunks:
        cid = c.metadata["chunk_id"]
        if cid in seen:
            continue
        seen.add(cid)
        unique.append(c)
    return unique


# --------------------------------------------------------------------------- #
# Incremental manifest                                                        #
# --------------------------------------------------------------------------- #
def _file_signature(path: Path) -> str:
    """Cheap change-detector: hash of size + mtime (no full re-read)."""
    st = path.stat()
    return hashlib.sha256(f"{st.st_size}:{int(st.st_mtime)}".encode()).hexdigest()[:16]


def _manifest_path(settings: Settings) -> Path:
    return settings.vector_dir.parent / MANIFEST_NAME


def _load_manifest(settings: Settings) -> dict[str, str]:
    p = _manifest_path(settings)
    if p.exists():
        try:
            return json.loads(p.read_text(encoding="utf-8"))
        except Exception:  # noqa: BLE001
            return {}
    return {}


def _save_manifest(settings: Settings, manifest: dict[str, str]) -> None:
    p = _manifest_path(settings)
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(manifest, indent=2), encoding="utf-8")


# --------------------------------------------------------------------------- #
# Indexing                                                                    #
# --------------------------------------------------------------------------- #
def build_index(
    settings: Settings | None = None, *, reset: bool = False
) -> dict[str, int]:
    """Load documents under ``data_dir``, chunk and embed the new/changed ones.

    Files whose size+mtime signature is unchanged since the last run are
    skipped (unless ``reset=True``). Returns a summary with counts.
    """
    settings = settings or get_settings()
    settings.ensure_dirs()

    files = discover_files(settings.data_dir)
    if not files:
        logger.warning("No supported documents found in %s", settings.data_dir)
        return {"files": 0, "chunks": 0, "skipped": 0}

    manifest = {} if reset else _load_manifest(settings)

    # Decide which files actually need (re)indexing.
    todo: list[Path] = []
    skipped = 0
    for f in files:
        sig = _file_signature(f)
        if manifest.get(str(f)) == sig:
            skipped += 1
        else:
            todo.append(f)

    if not todo:
        logger.info("Nothing to do — all %d files already indexed.", len(files))
        return {"files": len(files), "chunks": 0, "skipped": skipped}

    # Imported lazily so the pure load/chunk helpers don't require the vector
    # store (Chroma) dependency to be installed for unit tests.
    from .vectorstore import get_vectorstore

    store = get_vectorstore(settings, reset=reset)

    total_chunks = 0
    batch = 128
    for n, f in enumerate(todo, start=1):
        t0 = time.perf_counter()
        loaded = load_file(f)
        if not loaded:
            logger.warning("[%d/%d] %-40s — no extractable text", n, len(todo), f.name)
            manifest[str(f)] = _file_signature(f)  # don't retry empty files endlessly
            continue

        chunks = _dedupe_by_id(_chunk(loaded, settings))
        ids = [c.metadata["chunk_id"] for c in chunks]
        for i in range(0, len(chunks), batch):
            store.add_documents(chunks[i : i + batch], ids=ids[i : i + batch])

        manifest[str(f)] = _file_signature(f)
        total_chunks += len(chunks)
        dt = time.perf_counter() - t0
        logger.info(
            "[%d/%d] %-40s — %d sections → %d chunks (%.1fs)",
            n, len(todo), f.name, len(loaded), len(chunks), dt,
        )
        # Persist the manifest incrementally so an interrupted long run resumes.
        _save_manifest(settings, manifest)

    logger.info(
        "Indexed %d new/changed files (%d chunks); skipped %d unchanged.",
        len(todo), total_chunks, skipped,
    )
    return {"files": len(todo), "chunks": total_chunks, "skipped": skipped}
